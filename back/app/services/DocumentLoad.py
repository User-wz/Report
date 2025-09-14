# -*- coding:utf-8 -*-
# -*- coding=utf-8 -*-
# coding:utf-8
'''
project_name : win
file name : document_load
Author : Administrator
date : 2024/09/10  10:17
'''
import json
import os.path
import re
import fitz
from pathlib import Path
import PyPDF2
import docx
import ftfy
from langchain_community.document_loaders import BSHTMLLoader, JSONLoader
from langchain_community.document_loaders import PyPDFium2Loader
from langchain_community.document_loaders.csv_loader import CSVLoader


# 正则匹配优化文本读取内容
from app.services.resume import read_word_docx_in_order, print_ordered_content


def rm_newlines(text):
    text = re.sub(r'(?<=[^\.。:：])\n', ' ', text)
    text = re.sub(r'\t', ' ', text)
    text = re.sub(r'\n', ' ', text)
    text = re.sub(r'[\x00-\x1F\x7F-\x9F]', '', text)
    text = re.sub(r'[\ue5d2\ue5cf\ue5e5]', '', text)
    return text


def rm_cid(text):
    text = re.sub(r'\(cid:\d+\)', '', text)
    return text


def rm_hexadecimal(text):
    text = re.sub(r'[0-9A-Fa-f]{21,}', '', text)
    return text


def deal(text):
    text = rm_newlines(text)
    text = rm_cid(text)
    text = rm_hexadecimal(text)
    return text


class DocumentLoad():

    def remove_unicode(self, text):
        # return re.sub(r"\\[uU][0-9a-fA-F]+","",text)
        return ftfy.fix_text(text)

    def remove_(self, text):
        return "".join(re.findall(r"[^?•]+", text))

    def pdf_load(self, path_filename):
        with open(path_filename, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text_list = []
            def visitor_body(text, cm, tm, fontDict, fontSize):
                if fontDict and "/UniGB-UCS2-H" not in fontDict.values():
                    text_list.append(deal(text))
            # 遍历PDF的每一页并提取文本
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page.extract_text(visitor_text=visitor_body)

        return "".join(text_list)

    def pdf_loads(self,path_filename):
        from pdfminer.high_level import extract_text
        # doc = fitz.open(path_filename)
        # for page in doc:
        #     print(page.get_text())
        # from pdfminer.high_level import extract_text

        text = extract_text(path_filename)
        return text
    def json_load(self, path_filename):
        content = json.loads(Path(path_filename).read_text(encoding="utf-8"))
        content_str = ""
        if isinstance(content, dict):
            for key, value in content.items():
                if isinstance(value, list):
                    for v in value:
                        content_str += deal(str(v)) + "\n"
                else:
                    content_str += key + ":" + value
        if isinstance(content, list):
            for v in content:
                content_str += deal(str(v)) + "\n"
        return content_str


    def csv_load(self, path_filename):
        loader = CSVLoader(file_path=path_filename, encoding="utf-8")
        data = loader.load()
        content = ""
        for i in data:
            content += deal(str(i.page_content)) + "\n"
        return content

    def docx_load(self, path_filename):
        text_list = []
        doc = docx.Document(path_filename)
        for para in doc.paragraphs:
            if len(para.text.replace('\t', '').replace('\n', '')) > 0:
                # text_list.append(para.text.replace('\t', '').replace('\n', ''))
                text_list.append(deal(para.text))
        text = ' '.join(text_list)
        return text
    def docx_loads(self,path_filename):
        ordered_data = read_word_docx_in_order(path_filename)
        return print_ordered_content(ordered_data)

    def html_load(self, path_filename):
        loader = BSHTMLLoader(path_filename, open_encoding="utf-8")
        data = loader.load()
        print(data)

    def pptx_load(self, path_filename):
        from pptx import Presentation

        ppt = Presentation(path_filename)
        content = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    content += deal(shape.text_frame.text) + "\n"
        return content

    def txt_load(self, path_filename):
        text_list = []
        with open(path_filename, 'r', encoding='utf-8') as f:
            lines = f.readlines()
            for line in lines:
                if len(line.replace('\t', '').replace('\n', '')) > 0:
                    # text_list.append(line.replace('\t', '').replace('\n', ''))
                    text_list.append(deal(line))
        return ' '.join(text_list)

    # 替换文件内容
    def txt_replace(self, filePath, oldString, newString):
        print("正在处理文件内容 ... --> " + filePath)

        with open(filePath, mode="r", encoding="utf-8") as oldFile, open("%s.bak" % filePath, mode="w",
                                                                         encoding="utf-8") as newFile:
            oldFile.read()
            for line in oldFile:
                if oldString in line:
                    line = line.replace(oldString, newString)
                newFile.write(line)
        os.remove(filePath)
        os.rename("%s.bak" % filePath, filePath)
        print("文件内容处理完成！ --> " + filePath)
        return


if __name__ == '__main__':
    D = DocumentLoad()
    s = r"C:\Users\Administrator\Desktop\HR\20231212-韩漾简历-财务.pdf"
    with open("123.txt", 'w',encoding="utf-8")as fp:
        data=D.pdf_loads(s)
        fp.write(str(data))

